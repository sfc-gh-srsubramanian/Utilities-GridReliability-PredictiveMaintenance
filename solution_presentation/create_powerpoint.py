#!/usr/bin/env python3
"""
Create PowerPoint Presentation for Grid Reliability Solution
Upload the generated .pptx file to Google Slides

Uses Snowflake brand colors from SNOWFLAKE_COLOR_SCHEME.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Snowflake Brand Colors
COLORS = {
    'snowflake_blue': RGBColor(41, 181, 232),      # #29B5E8
    'mid_blue': RGBColor(17, 86, 127),            # #11567F
    'midnight': RGBColor(0, 0, 0),                 # #000000
    'star_blue': RGBColor(113, 211, 220),          # #71D3DC
    'valencia_orange': RGBColor(255, 159, 54),     # #FF9F36
    'purple_moon': RGBColor(125, 68, 207),         # #7D44CF
    'first_light': RGBColor(212, 91, 144),         # #D45B90
    'windy_city': RGBColor(138, 153, 158),         # #8A999E
    'iceberg': RGBColor(0, 53, 69),                # #003545
    'winter': RGBColor(36, 50, 61),                # #24323D
    'white': RGBColor(255, 255, 255),
    'light_gray': RGBColor(242, 242, 242),
}

def create_title_slide(prs, title, subtitle, bg_color='winter'):
    """Create a title slide with Snowflake branding"""
    
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS[bg_color]
    
    # Add title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Format title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['snowflake_blue']
    title_para.font.name = 'Arial'
    
    # Add subtitle
    left = Inches(0.5)
    top = Inches(4.2)
    width = Inches(9)
    height = Inches(1.8)
    
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    subtitle_frame.word_wrap = True
    
    # Format subtitle
    for para in subtitle_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = COLORS['white']
        para.font.name = 'Arial'
        para.space_after = Pt(10)
    
    return slide

def create_content_slide(prs, title, subtitle, sections):
    """Create a standard content slide with blue header"""
    
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['white']
    
    # Add blue header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['snowflake_blue']
    header.line.fill.background()
    
    # Add title on header
    title_box = header.text_frame
    title_box.text = title
    title_box.margin_left = Inches(0.2)
    title_box.margin_top = Inches(0.1)
    title_box.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    title_para = title_box.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.font.name = 'Arial'
    
    # Add subtitle if present
    current_top = 1.0
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(current_top),
            Inches(9.4), Inches(0.4)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.italic = True
        subtitle_para.font.color.rgb = COLORS['mid_blue']
        subtitle_para.font.name = 'Arial'
        
        current_top = 1.5
    
    # Add content sections
    content_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(current_top),
        Inches(9.4), Inches(7.5 - current_top)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for section in sections:
        # Section heading
        p = content_frame.add_paragraph()
        p.text = section['heading']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['midnight']
        p.font.name = 'Arial'
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        
        # Bullets
        for bullet in section['bullets']:
            p = content_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['midnight']
            p.font.name = 'Arial'
            p.space_after = Pt(3)
    
    # Remove the default first paragraph
    if content_frame.paragraphs[0].text == '':
        content_frame.paragraphs[0]._element.getparent().remove(content_frame.paragraphs[0]._element)
    
    return slide

def create_impact_slide(prs, title, subtitle, stats, sections, footer=''):
    """Create impact slide with highlighted statistics"""
    
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Light gray background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['light_gray']
    
    # Blue header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.9)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['snowflake_blue']
    header.line.fill.background()
    
    # Title
    title_box = header.text_frame
    title_box.text = title
    title_box.margin_left = Inches(0.2)
    title_box.margin_top = Inches(0.1)
    title_box.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    title_para = title_box.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.font.name = 'Arial'
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(1.0),
            Inches(9.4), Inches(0.4)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.italic = True
        subtitle_para.font.color.rgb = COLORS['mid_blue']
        subtitle_para.font.name = 'Arial'
    
    # Stats boxes
    if stats:
        stat_width = 1.7
        stat_height = 1.2
        stat_spacing = 0.1
        start_x = 0.5
        start_y = 1.8
        
        for i, stat in enumerate(stats):
            x = start_x + (i % 5) * (stat_width + stat_spacing)
            y = start_y + (i // 5) * (stat_height + stat_spacing)
            
            # Stat box
            stat_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(stat_width), Inches(stat_height)
            )
            stat_box.fill.solid()
            stat_box.fill.fore_color.rgb = COLORS['white']
            stat_box.line.color.rgb = COLORS['snowflake_blue']
            stat_box.line.width = Pt(2)
            
            # Icon + Stat + Description
            stat_frame = stat_box.text_frame
            stat_frame.text = f"{stat['icon']}\n{stat['stat']}\n{stat['desc']}"
            stat_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Format
            for idx, para in enumerate(stat_frame.paragraphs):
                para.alignment = PP_ALIGN.CENTER
                if idx == 0:  # Icon
                    para.font.size = Pt(24)
                elif idx == 1:  # Stat
                    para.font.size = Pt(18)
                    para.font.bold = True
                    para.font.color.rgb = COLORS['snowflake_blue']
                else:  # Description
                    para.font.size = Pt(11)
                    para.font.color.rgb = COLORS['midnight']
                para.font.name = 'Arial'
    
    # Additional sections
    if sections:
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.0),
            Inches(9), Inches(2.5)
        )
        content_frame = content_box.text_frame
        
        for section in sections:
            p = content_frame.add_paragraph()
            p.text = section['heading']
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = COLORS['midnight']
            p.font.name = 'Arial'
            p.space_after = Pt(6)
            
            for bullet in section['bullets']:
                p = content_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS['midnight']
                p.font.name = 'Arial'
        
        if content_frame.paragraphs[0].text == '':
            content_frame.paragraphs[0]._element.getparent().remove(content_frame.paragraphs[0]._element)
    
    # Footer
    if footer:
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8),
            Inches(9), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = footer
        
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(10)
        footer_para.font.italic = True
        footer_para.font.color.rgb = COLORS['windy_city']
        footer_para.font.name = 'Arial'
    
    return slide

def main():
    """Generate complete presentation"""
    
    print("\n" + "="*70)
    print("  GRID RELIABILITY POWERPOINT GENERATOR")
    print("  Using Snowflake Brand Colors")
    print("="*70 + "\n")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("📊 Generating slides...\n")
    
    # Slide 1: Title
    print("  ✅ Slide 1: Title")
    create_title_slide(
        prs,
        "Grid Reliability & Predictive Maintenance",
        "AI-Powered Asset Intelligence on Snowflake\n\nTransforming Utility Operations Through Intelligence-Driven Maintenance",
        'winter'
    )
    
    # Slide 2: The Challenge
    print("  ✅ Slide 2: The Utility Challenge")
    create_content_slide(
        prs,
        "The Utility Challenge",
        "Modern Grid Operations Face Critical Pressures",
        [
            {
                'heading': 'Aging Infrastructure Crisis',
                'bullets': [
                    '40% of transformers & circuit breakers >20 years old',
                    'Traditional 25-year design life under stress',
                    'Delayed failures = customer impact + penalties'
                ]
            },
            {
                'heading': 'Reactive Maintenance is Failing',
                'bullets': [
                    '60-70% of failures occur DESPITE calendar-based maintenance',
                    'Emergency replacements cost 3-5x more than planned',
                    'Average transformer failure: $385K + 4.2 hour outage + 8,500 customers'
                ]
            },
            {
                'heading': 'Data Trapped in Silos',
                'bullets': [
                    'OT sensor data in SCADA (Supervisory Control and Data Acquisition)',
                    'IT asset data in separate enterprise systems',
                    'Maintenance logs & manuals not analyzed',
                    'Visual inspections (drone, thermal) underutilized'
                ]
            }
        ]
    )
    
    # Slide 3: Regulatory Pressures
    print("  ✅ Slide 3: Regulatory & Climate Pressures")
    create_content_slide(
        prs,
        "Regulatory & Climate Pressures",
        "Utilities Under Unprecedented Pressure",
        [
            {
                'heading': 'Regulatory Scrutiny',
                'bullets': [
                    'State commissions closely monitor SAIDI/SAIFI metrics',
                    '  → SAIDI: System Average Interruption Duration Index',
                    '  → SAIFI: System Average Interruption Frequency Index',
                    'Financial penalties for poor reliability performance',
                    'Pressure to justify rate cases with improvements'
                ]
            },
            {
                'heading': 'Climate & Load Growth',
                'bullets': [
                    'Extreme weather increasing thermal stress',
                    'EV adoption driving unprecedented load',
                    'Aging infrastructure meets growing demand',
                    'Grid modernization imperative'
                ]
            },
            {
                'heading': 'Industry Benchmarks (EIA 2023)',
                'bullets': [
                    'U.S. Average SAIDI: 118.4 minutes/customer/year',
                    'U.S. Average SAIFI: 0.999 interruptions/customer/year'
                ]
            }
        ]
    )
    
    # Slide 4: Business Opportunity
    print("  ✅ Slide 4: The Business Opportunity")
    create_impact_slide(
        prs,
        "The Business Opportunity",
        "AI-Powered Predictive Maintenance Delivers Results",
        [
            {'icon': '⚡', 'stat': '70% reduction', 'desc': 'in unplanned outages'},
            {'icon': '💰', 'stat': '$25M+ annual', 'desc': 'cost avoidance'},
            {'icon': '📈', 'stat': '15-25%', 'desc': 'SAIDI/SAIFI improvement'},
            {'icon': '🔧', 'stat': '40% reduction', 'desc': 'in maintenance costs'},
            {'icon': '⏱️', 'stat': '5-7 years', 'desc': 'asset life extension'}
        ],
        [
            {
                'heading': 'ROI Example (Mid-Sized Utility)',
                'bullets': [
                    'Investment: $2-3M over 18-24 months',
                    'Annual Benefit: $15-35M',
                    'Net ROI: 500-1,500% over 3 years',
                    'Payback: 2-6 months'
                ]
            }
        ],
        'Source: DOE Grid Modernization Reports, EPRI Asset Management Studies, IEEE Power & Energy Society'
    )
    
    # Slide 5: Solution Overview
    print("  ✅ Slide 5: Our Solution - Platform Overview")
    create_content_slide(
        prs,
        "Our Solution - Platform Overview",
        "Comprehensive AI-Powered Predictive Maintenance",
        [
            {
                'heading': 'Unified Data Foundation',
                'bullets': [
                    '360° asset health monitoring across IT + OT systems',
                    'Single source of truth on Snowflake AI Data Cloud'
                ]
            },
            {
                'heading': 'Intelligent Predictions',
                'bullets': [
                    'ML models predict failures 14-30 days in advance',
                    'Anomaly detection & Remaining Useful Life (RUL) estimation',
                    'Real-time risk scoring with confidence levels'
                ]
            },
            {
                'heading': 'Unstructured Intelligence',
                'bullets': [
                    'Analyzes maintenance logs, technical manuals',
                    'Visual inspection data (drone, thermal, LiDAR)',
                    'Computer Vision detection integration'
                ]
            },
            {
                'heading': 'Natural Language Access',
                'bullets': [
                    'Conversational analytics via Snowflake Intelligence Agents',
                    'No SQL required for operators'
                ]
            }
        ]
    )
    
    # Add a few more key slides
    print("  ✅ Slide 6-10: Additional slides (use the pattern above to add more)")
    
    # Save
    filename = "Grid_Reliability_Snowflake_Presentation.pptx"
    prs.save(filename)
    
    print(f"\n{'='*70}")
    print(f"  ✅ COMPLETE!")
    print(f"{'='*70}")
    print(f"\n📁 File saved: {filename}")
    print(f"\n📤 Next steps:")
    print(f"   1. Open {filename} to review")
    print(f"   2. Upload to Google Drive")
    print(f"   3. Right-click → Open with → Google Slides")
    print(f"   4. Add remaining slides using the patterns in the script")
    print(f"   5. Add images from solution_presentation/images/\n")

if __name__ == '__main__':
    main()
